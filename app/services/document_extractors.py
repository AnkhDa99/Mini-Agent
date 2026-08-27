import csv
import json
import os
from dataclasses import dataclass
from typing import Iterable
from xml.etree import ElementTree

from bs4 import BeautifulSoup
from pypdf import PdfReader
from docx import Document as DocxDocument
from pptx import Presentation
from openpyxl import load_workbook


@dataclass
class ExtractedBlock:
    content: str
    source_type: str
    page_no: int | None = None
    sheet_name: str | None = None
    slide_no: int | None = None


class DocumentExtractor:
    """
    统一文档解析入口：
    不同文件类型 -> ExtractedBlock 列表

    注意：
    这里先解决“能解析大量项目资料”的问题。
    后续可以继续扩展 OCR、Apache Tika、LibreOffice 转换等能力。
    """

    TEXT_EXTS = {
        ".txt", ".md", ".markdown",
        ".py", ".java", ".js", ".ts", ".vue",
        ".html", ".css", ".sql", ".json",
        ".yaml", ".yml", ".xml", ".log",
    }

    def extract(self, local_path: str) -> Iterable[ExtractedBlock]:
        ext = os.path.splitext(local_path)[1].lower()

        if ext in {".txt", ".md", ".markdown", ".py", ".java", ".js", ".ts", ".vue", ".css", ".sql", ".yaml", ".yml", ".log"}:
            yield from self._extract_plain_text(local_path, source_type=ext.lstrip("."))
            return

        if ext == ".json":
            yield from self._extract_json(local_path)
            return

        if ext == ".xml":
            yield from self._extract_xml(local_path)
            return

        if ext in {".html", ".htm"}:
            yield from self._extract_html(local_path)
            return

        if ext == ".pdf":
            yield from self._extract_pdf(local_path)
            return

        if ext == ".docx":
            yield from self._extract_docx(local_path)
            return

        if ext == ".pptx":
            yield from self._extract_pptx(local_path)
            return

        if ext in {".xlsx", ".xlsm"}:
            yield from self._extract_xlsx(local_path)
            return

        if ext == ".csv":
            yield from self._extract_csv(local_path)
            return

        raise ValueError(f"暂不支持该文件类型：{ext}")

    def _read_text_with_fallback(self, local_path: str) -> str:
        for encoding in ["utf-8", "utf-8-sig", "gbk", "gb18030"]:
            try:
                with open(local_path, "r", encoding=encoding) as f:
                    return f.read()
            except UnicodeDecodeError:
                continue

        raise ValueError("文件编码无法识别")

    def _extract_plain_text(self, local_path: str, source_type: str) -> Iterable[ExtractedBlock]:
        text = self._read_text_with_fallback(local_path)
        if text.strip():
            yield ExtractedBlock(
                content=text,
                source_type=source_type,
            )

    def _extract_json(self, local_path: str) -> Iterable[ExtractedBlock]:
        raw = self._read_text_with_fallback(local_path)

        try:
            obj = json.loads(raw)
            text = json.dumps(obj, ensure_ascii=False, indent=2)
        except Exception:
            text = raw

        if text.strip():
            yield ExtractedBlock(
                content=text,
                source_type="json",
            )

    def _extract_xml(self, local_path: str) -> Iterable[ExtractedBlock]:
        raw = self._read_text_with_fallback(local_path)

        try:
            root = ElementTree.fromstring(raw)
            texts = []

            def walk(node, depth=0):
                tag = node.tag
                text = (node.text or "").strip()
                if text:
                    texts.append(f"{'  ' * depth}<{tag}> {text}")
                for child in list(node):
                    walk(child, depth + 1)

            walk(root)
            text = "\n".join(texts) if texts else raw

        except Exception:
            text = raw

        if text.strip():
            yield ExtractedBlock(
                content=text,
                source_type="xml",
            )

    def _extract_html(self, local_path: str) -> Iterable[ExtractedBlock]:
        raw = self._read_text_with_fallback(local_path)
        soup = BeautifulSoup(raw, "html.parser")

        for tag in soup(["script", "style"]):
            tag.extract()

        text = soup.get_text(separator="\n")
        lines = [line.strip() for line in text.splitlines() if line.strip()]
        text = "\n".join(lines)

        if text.strip():
            yield ExtractedBlock(
                content=text,
                source_type="html",
            )

    def _extract_pdf(self, local_path: str) -> Iterable[ExtractedBlock]:
        reader = PdfReader(local_path)

        for page_index, page in enumerate(reader.pages):
            text = page.extract_text() or ""
            text = text.strip()

            if text:
                yield ExtractedBlock(
                    content=text,
                    source_type="pdf",
                    page_no=page_index + 1,
                )

    def _extract_docx(self, local_path: str) -> Iterable[ExtractedBlock]:
        doc = DocxDocument(local_path)
        parts = []

        for p in doc.paragraphs:
            text = p.text.strip()
            if text:
                parts.append(text)

        for table_index, table in enumerate(doc.tables):
            rows = []
            for row in table.rows:
                cells = [cell.text.strip().replace("\n", " ") for cell in row.cells]
                rows.append(" | ".join(cells))
            if rows:
                parts.append(f"[TABLE {table_index + 1}]\n" + "\n".join(rows))

        text = "\n\n".join(parts)

        if text.strip():
            yield ExtractedBlock(
                content=text,
                source_type="docx",
            )

    def _extract_pptx(self, local_path: str) -> Iterable[ExtractedBlock]:
        prs = Presentation(local_path)

        for slide_index, slide in enumerate(prs.slides):
            texts = []

            for shape in slide.shapes:
                if hasattr(shape, "text"):
                    text = shape.text.strip()
                    if text:
                        texts.append(text)

            content = "\n".join(texts).strip()

            if content:
                yield ExtractedBlock(
                    content=content,
                    source_type="pptx",
                    slide_no=slide_index + 1,
                )

    def _extract_xlsx(self, local_path: str) -> Iterable[ExtractedBlock]:
        wb = load_workbook(local_path, read_only=True, data_only=True)

        for ws in wb.worksheets:
            lines = []
            for row in ws.iter_rows():
                values = []
                for cell in row:
                    value = cell.value
                    if value is None:
                        values.append("")
                    else:
                        values.append(str(value).strip())

                if any(values):
                    lines.append(" | ".join(values))

            content = "\n".join(lines).strip()

            if content:
                yield ExtractedBlock(
                    content=content,
                    source_type="xlsx",
                    sheet_name=ws.title,
                )

    def _extract_csv(self, local_path: str) -> Iterable[ExtractedBlock]:
        for encoding in ["utf-8", "utf-8-sig", "gbk", "gb18030"]:
            try:
                lines = []
                with open(local_path, "r", encoding=encoding, newline="") as f:
                    reader = csv.reader(f)
                    for row in reader:
                        lines.append(" | ".join(row))

                content = "\n".join(lines).strip()

                if content:
                    yield ExtractedBlock(
                        content=content,
                        source_type="csv",
                    )
                return

            except UnicodeDecodeError:
                continue

        raise ValueError("CSV 文件编码无法识别")