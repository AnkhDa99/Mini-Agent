"""
文档生成服务：基于检索到的项目上下文 + LLM 内容生成，产出 PPT / Word / PDF / Excel / Mermaid 图表。

生成后的文件保存到 data/output/ 目录，返回文件路径供前端下载。
"""
from __future__ import annotations

import logging
import os
import tempfile
from dataclasses import dataclass
from pathlib import Path

logger = logging.getLogger(__name__)

_OUTPUT_DIR: Path | None = None


def _ensure_output_dir() -> Path:
    """确保输出目录存在。"""
    global _OUTPUT_DIR
    if _OUTPUT_DIR is None:
        from app.core.config import settings
        _OUTPUT_DIR = Path(settings.document_output_dir)
    _OUTPUT_DIR.mkdir(parents=True, exist_ok=True)
    return _OUTPUT_DIR


@dataclass
class DocResult:
    """文档生成结果。"""
    filepath: str
    filename: str
    mime_type: str
    content: str = ""  # 附加内容（如图表 Mermaid 代码）
    error: str = ""


def _generate_title(llm_client, user_query: str, doc_type: str, context: str = "") -> str:
    """用 LLM 生成一个简洁概括的文档标题（不超过 30 字）。"""
    prompt = f"""根据用户需求和文档上下文，为一份{doc_type}生成一个简洁的标题。

用户需求: {user_query}
{f"参考上下文: {context[:2000]}" if context else ""}

要求:
- 不超过 30 个字
- 概括文档的核心主题
- 不要包含引号、斜杠等特殊字符
- 只输出标题，不要任何解释

标题:"""
    try:
        title = llm_client.chat([{"role": "user", "content": prompt}], temperature=0.3)
        title = title.strip().strip("\"'").replace("\n", " ")
        if len(title) > 50:
            title = title[:50]
        return title or user_query[:40]
    except Exception:
        logger.exception("Title generation failed")
        return user_query[:40]


def _build_content_prompt(user_query: str, title: str, doc_type: str,
                          context: str, extra: str = "") -> str:
    """构建 LLM 内容生成的 prompt。"""
    base = f"""基于以下项目文档片段，生成一份{doc_type}的内容。

标题: {title}
用户需求: {user_query}
{extra}
## 参考文档上下文
{context}

## 输出要求
- 使用 Markdown 格式
- 结构清晰，有层级标题
- 严格基于文档上下文提取和整理内容，不要编造
- 文档未覆盖的部分，标注 [通用知识]
- 严禁出现"参考XX文档"、"请参考"、"详见"等引用性语句，必须直接呈现内容本身"""
    return base + "\n\n请生成完整的文档内容："


def generate_markdown(llm_client, user_query: str, title: str,
                      context: str) -> DocResult:
    """生成 Markdown (.md) 文件。"""
    prompt = _build_content_prompt(user_query, title, "Markdown 文档", context)

    try:
        raw = llm_client.chat([{"role": "user", "content": prompt}], temperature=0.5)
    except Exception:
        logger.exception("LLM生成Markdown内容失败")
        raw = f"# {title}\n\n基于检索上下文生成的内容。\n\n文档库未找到相关内容。[通用知识]"

    out_dir = _ensure_output_dir()
    safe_name = _safe_filename(title)
    filepath = out_dir / f"{safe_name}.md"
    filepath.write_text(raw, encoding="utf-8")

    logger.info("Markdown generated: %s (%d chars)", filepath, len(raw))
    return DocResult(
        filepath=str(filepath),
        filename=f"{safe_name}.md",
        mime_type="text/markdown",
    )


def generate_ppt(llm_client, user_query: str, title: str, outline: str,
                 context: str) -> DocResult:
    """生成 PPT 文件。

    Args:
        llm_client: LLM 客户端
        user_query: 用户原始问题
        title: PPT 标题
        outline: 大纲（逗号分隔的页标题，若为空则由 LLM 生成）
        context: 检索到的文档上下文
    """
    try:
        from pptx import Presentation
        from pptx.util import Inches, Pt
        from pptx.enum.text import PP_ALIGN
    except ImportError:
        return DocResult(filepath="", filename="", mime_type="",
                         error="python-pptx 未安装，无法生成 PPT")

    # 用 LLM 生成大纲（如果未提供有效大纲）
    slides_outline = [s.strip() for s in outline.split(",") if s.strip() and s.strip() != title]
    if len(slides_outline) < 3:
        outline_prompt = f"""基于以下内容，为一份 PPT 生成 4-6 页的幻灯片大纲。

用户需求: {user_query}
参考上下文: {context[:3000]}

要求: 每行一个页标题，不要编号，不要解释。每页标题 10 字以内。
大纲:"""
        try:
            raw_outline = llm_client.chat([{"role": "user", "content": outline_prompt}], temperature=0.3)
            slides_outline = [s.strip().lstrip("-•· ").strip() for s in raw_outline.split("\n") if s.strip()]
            if len(slides_outline) < 3:
                slides_outline = ["项目背景与目标", "核心概念与架构", "关键技术要点", "实现方案分析", "总结与建议"]
        except Exception:
            slides_outline = ["项目背景与目标", "核心概念与架构", "关键技术要点", "实现方案分析", "总结与建议"]

    prompt = f"""基于以下项目文档片段，生成一份 PPT 演示文稿的内容。

标题: {title}
用户需求: {user_query}
幻灯片大纲: {', '.join(slides_outline)}

## 参考文档上下文
{context}

## 输出格式（严格遵守）
为每一页幻灯片输出:
## <页标题>
- <具体要点1——从文档中提取的事实/数据/观点>
- <具体要点2——从文档中提取的事实/数据/观点>
- <具体要点3>
(每页 3-5 条要点，每条要点 15-30 字)

## 输出规则（非常重要）
- 每条要点必须是具体的、有信息量的陈述或事实
- 严禁出现"参考XX文档"、"详见"、"请参考"、"相关内容"等占位性表述
- 必须从文档上下文中提取实际内容填充每一页
- 文档未覆盖的部分，标注 [通用知识]
- 只输出幻灯片内容，不要额外解释"""

    try:
        raw = llm_client.chat([{"role": "user", "content": prompt}], temperature=0.5)
        slides_content = _parse_slides(raw, slides_outline)
    except Exception:
        logger.exception("LLM生成PPT内容失败")
        slides_content = {s: [f"• [通用知识] {s}相关内容"] for s in slides_outline}

    # 创建 PPT
    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)

    # 封面
    slide = prs.slides.add_slide(prs.slide_layouts[6])  # 空白版式
    _add_text_box(slide, Inches(1), Inches(2.5), Inches(11.3), Inches(1.5),
                  title, font_size=Pt(36), bold=True, alignment=PP_ALIGN.CENTER)
    _add_text_box(slide, Inches(1), Inches(4.2), Inches(11.3), Inches(0.8),
                  f"Mini Agent 自动生成", font_size=Pt(16), alignment=PP_ALIGN.CENTER)

    # 内容页
    for s_title in slides_outline:
        slide = prs.slides.add_slide(prs.slide_layouts[6])
        _add_text_box(slide, Inches(0.8), Inches(0.4), Inches(11.7), Inches(0.9),
                      s_title, font_size=Pt(28), bold=True)
        content = slides_content.get(s_title, [f"• {s_title}"])
        _add_text_box(slide, Inches(1.2), Inches(1.6), Inches(10.9), Inches(5.2),
                      "\n".join(content), font_size=Pt(18))

    # 保存
    out_dir = _ensure_output_dir()
    safe_name = _safe_filename(title)
    filepath = out_dir / f"{safe_name}.pptx"
    prs.save(str(filepath))

    logger.info("PPT generated: %s (%d slides)", filepath, len(slides_outline) + 1)
    return DocResult(
        filepath=str(filepath),
        filename=f"{safe_name}.pptx",
        mime_type="application/vnd.openxmlformats-officedocument.presentationml.presentation",
    )


def generate_word(llm_client, user_query: str, title: str, content_type: str,
                  context: str) -> DocResult:
    """生成 Word 文档。"""
    try:
        from docx import Document
        from docx.shared import Pt, Inches, RGBColor
        from docx.enum.text import WD_ALIGN_PARAGRAPH
    except ImportError:
        return DocResult(filepath="", filename="", mime_type="",
                         error="python-docx 未安装，无法生成 Word")

    prompt = _build_content_prompt(
        user_query, title, f"Word {content_type}", context,
    )

    try:
        raw = llm_client.chat([{"role": "user", "content": prompt}], temperature=0.5)
    except Exception:
        logger.exception("LLM生成Word内容失败")
        raw = f"# {title}\n\n基于检索上下文生成的内容。\n\n文档库未找到相关内容。[通用知识]"

    doc = Document()

    # 标题
    h = doc.add_heading(title, level=0)
    h.alignment = WD_ALIGN_PARAGRAPH.CENTER

    doc.add_paragraph(f"文档类型: {content_type}")
    doc.add_paragraph(f"生成方式: Mini Agent 自动生成")
    doc.add_paragraph("")

    # 解析 Markdown 转为 Word 段落
    in_code_block = False
    code_lines = []
    for line in raw.split("\n"):
        stripped = line.strip()
        if stripped.startswith("```"):
            if in_code_block:
                doc.add_paragraph("\n".join(code_lines), style="Code")
                code_lines = []
                in_code_block = False
            else:
                in_code_block = True
            continue
        if in_code_block:
            code_lines.append(line)
            continue
        if stripped.startswith("# ") and not stripped.startswith("## "):
            doc.add_heading(stripped[2:], level=1)
        elif stripped.startswith("## "):
            doc.add_heading(stripped[3:], level=2)
        elif stripped.startswith("### "):
            doc.add_heading(stripped[4:], level=3)
        elif stripped.startswith("- ") or stripped.startswith("* "):
            doc.add_paragraph(stripped, style="List Bullet")
        elif stripped:
            doc.add_paragraph(stripped)

    out_dir = _ensure_output_dir()
    safe_name = _safe_filename(title)
    filepath = out_dir / f"{safe_name}.docx"
    doc.save(str(filepath))

    logger.info("Word generated: %s", filepath)
    return DocResult(
        filepath=str(filepath),
        filename=f"{safe_name}.docx",
        mime_type="application/vnd.openxmlformats-officedocument.wordprocessingml.document",
    )


def generate_pdf(llm_client, user_query: str, title: str,
                 context: str) -> DocResult:
    """生成 PDF 文档。

    优先使用 fpdf2，其次 reportlab，都不可用则尝试用 python-docx + 系统打印。
    """
    prompt = _build_content_prompt(user_query, title, "PDF 文档", context)

    try:
        raw = llm_client.chat([{"role": "user", "content": prompt}], temperature=0.5)
    except Exception:
        logger.exception("LLM生成PDF内容失败")
        raw = f"# {title}\n\n基于检索上下文生成的内容。"

    out_dir = _ensure_output_dir()
    safe_name = _safe_filename(title)

    # 尝试 fpdf2
    try:
        from fpdf import FPDF
        filepath = _generate_pdf_fpdf(raw, title, out_dir, safe_name)
        if filepath:
            return DocResult(
                filepath=str(filepath),
                filename=f"{safe_name}.pdf",
                mime_type="application/pdf",
            )
    except ImportError:
        logger.info("fpdf2 not available, trying reportlab...")

    # 尝试 reportlab
    try:
        from reportlab.lib.pagesizes import A4
        from reportlab.platypus import SimpleDocTemplate, Paragraph, Spacer
        from reportlab.lib.styles import getSampleStyleSheet
        filepath = out_dir / f"{safe_name}.pdf"
        doc = SimpleDocTemplate(str(filepath), pagesize=A4)
        styles = getSampleStyleSheet()
        story = [Paragraph(f"<b>{title}</b>", styles["Title"]), Spacer(1, 12)]
        for line in raw.split("\n"):
            stripped = line.strip()
            if stripped.startswith("# "):
                story.append(Paragraph(f"<b>{stripped[2:]}</b>", styles["Heading1"]))
            elif stripped.startswith("## "):
                story.append(Paragraph(f"<b>{stripped[3:]}</b>", styles["Heading2"]))
            elif stripped:
                story.append(Paragraph(stripped, styles["Normal"]))
        doc.build(story)
        logger.info("PDF generated (reportlab): %s", filepath)
        return DocResult(
            filepath=str(filepath),
            filename=f"{safe_name}.pdf",
            mime_type="application/pdf",
        )
    except ImportError:
        logger.info("reportlab not available either")

    return DocResult(filepath="", filename="", mime_type="",
                     error="PDF 生成需要 fpdf2 或 reportlab，请安装其中之一: pip install fpdf2")


def _generate_pdf_fpdf(raw: str, title: str, out_dir: Path, safe_name: str) -> Path | None:
    """用 fpdf2 生成 PDF。支持 CJK 需额外字体。"""
    from fpdf import FPDF

    # 尝试找到中文字体
    cjk_font_path = _find_cjk_font()
    pdf = FPDF()
    pdf.set_auto_page_break(auto=True, margin=15)
    pdf.add_page()

    if cjk_font_path:
        pdf.add_font("CJK", "", str(cjk_font_path), uni=True)
        pdf.add_font("CJK", "B", str(cjk_font_path), uni=True)
        title_font = ("CJK", "B", 16)
        body_font = ("CJK", "", 11)
    else:
        title_font = ("Helvetica", "B", 16)
        body_font = ("Helvetica", "", 11)
        logger.warning("No CJK font found, PDF may not render Chinese correctly")

    pdf.set_font(*title_font)
    pdf.cell(0, 10, title, ln=True, align="C")
    pdf.ln(6)

    pdf.set_font(*body_font)
    for line in raw.split("\n"):
        stripped = line.strip()
        if stripped.startswith("```"):
            continue
        if stripped.startswith("#"):
            pdf.set_font(*(title_font[0], title_font[1], 12))
            pdf.ln(4)
            pdf.cell(0, 8, stripped.lstrip("# "), ln=True)
            pdf.set_font(*body_font)
        elif stripped:
            pdf.multi_cell(0, 6, stripped)
            pdf.ln(1)

    filepath = out_dir / f"{safe_name}.pdf"
    pdf.output(str(filepath))
    logger.info("PDF generated (fpdf2): %s", filepath)
    return filepath


def generate_excel(llm_client, user_query: str, title: str, sheet_type: str,
                   context: str) -> DocResult:
    """生成 Excel 表格。"""
    try:
        import openpyxl
        from openpyxl.styles import Font, Alignment, PatternFill
    except ImportError:
        return DocResult(filepath="", filename="", mime_type="",
                         error="openpyxl 未安装，无法生成 Excel")

    prompt = _build_content_prompt(
        user_query, title, f"Excel {sheet_type}", context,
        extra="请以 Markdown 表格格式输出数据，可以有多个表格，每个表格前用 ## 表格名 标注。"
    )

    try:
        raw = llm_client.chat([{"role": "user", "content": prompt}], temperature=0.5)
    except Exception:
        logger.exception("LLM生成Excel内容失败")
        raw = f"## {sheet_type}\n\n无法生成内容。请稍后重试。"

    wb = openpyxl.Workbook()
    # 移除默认 sheet
    wb.remove(wb.active)

    # 解析 Markdown 表格
    tables = _parse_markdown_tables(raw)

    if not tables:
        ws = wb.create_sheet(title=title[:31])
        ws["A1"] = title
        ws["A1"].font = Font(size=14, bold=True)
        ws["A2"] = "未能解析到表格数据，请检查检索结果。"
    else:
        for t_name, headers, rows in tables:
            sheet_title = (t_name or sheet_type)[:31]
            ws = wb.create_sheet(title=sheet_title)

            header_font = Font(bold=True, color="FFFFFF", size=11)
            header_fill = PatternFill(start_color="4472C4", end_color="4472C4", fill_type="solid")
            header_align = Alignment(horizontal="center", vertical="center")

            for ci, h in enumerate(headers, 1):
                cell = ws.cell(row=1, column=ci, value=h)
                cell.font = header_font
                cell.fill = header_fill
                cell.alignment = header_align

            for ri, row in enumerate(rows, 2):
                for ci, val in enumerate(row, 1):
                    ws.cell(row=ri, column=ci, value=val)

            # 自动调整列宽
            for ci, h in enumerate(headers, 1):
                max_len = len(h) * 2
                for ri in range(2, len(rows) + 2):
                    cell_val = str(ws.cell(row=ri, column=ci).value or "")
                    max_len = max(max_len, len(cell_val) * 1.2)
                ws.column_dimensions[openpyxl.utils.get_column_letter(ci)].width = min(max_len, 50)

    out_dir = _ensure_output_dir()
    safe_name = _safe_filename(title)
    filepath = out_dir / f"{safe_name}.xlsx"
    wb.save(str(filepath))

    logger.info("Excel generated: %s (%d sheets)", filepath, len(tables))
    return DocResult(
        filepath=str(filepath),
        filename=f"{safe_name}.xlsx",
        mime_type="application/vnd.openxmlformats-officedocument.spreadsheetml.sheet",
    )


def generate_diagram(llm_client, user_query: str, diagram_type: str,
                     description: str, context: str) -> DocResult:
    """生成 Mermaid 图表代码。

    Args:
        llm_client: LLM 客户端
        user_query: 用户原始问题
        diagram_type: architecture / flowchart / sequence / class
        description: 图表内容描述
        context: 检索到的文档上下文
    """
    type_hints = {
        "architecture": "系统架构图 (graph TD) — 展示系统组件及其数据流/依赖关系",
        "flowchart": "业务流程图 (flowchart TD) — 展示决策分支和处理步骤",
        "sequence": "时序图 (sequenceDiagram) — 展示组件间交互消息顺序",
        "class": "类图 (classDiagram) — 展示实体及其属性和关系",
    }
    hint = type_hints.get(diagram_type, "使用 Mermaid graph TD 语法")

    prompt = f"""基于以下项目文档上下文，生成一份专业的{hint}的 Mermaid 代码。

图表主题: {description}
用户需求: {user_query}

## 参考文档上下文
{context}

## Mermaid 代码要求
1. 以以下主题初始化代码开头:
```
%%{{init: {{'theme': 'base', 'themeVariables': {{ 'primaryColor': '#e8f4fd', 'primaryTextColor': '#1a365d', 'primaryBorderColor': '#3182ce', 'lineColor': '#718096', 'secondaryColor': '#f0fff4', 'tertiaryColor': '#faf5ff' }}}}}}%%
```
2. 节点使用中文标注，文字简洁（每个节点不超过 12 字）
3. 使用 subgraph 对相关模块分组
4. 合理使用不同形状区分类型: 方框[流程]、圆角(起止)、菱形{{判断}}
5. 图表应反映文档中的实际组件/流程，不要虚构
6. 信息不足时合理推测但标注 [推测]
7. 节点数控制在 8-20 个，层级不超过 4 层
8. 只输出 Mermaid 代码，不要任何解释或 markdown 包裹

直接输出 Mermaid 代码（不含 ```mermaid 标记）:"""

    try:
        mermaid_code = llm_client.chat([{"role": "user", "content": prompt}], temperature=0.3)
        mermaid_code = mermaid_code.strip()
        # 清理可能的 markdown 包裹
        if mermaid_code.startswith("```mermaid"):
            mermaid_code = mermaid_code[len("```mermaid"):].strip()
        if mermaid_code.startswith("```"):
            mermaid_code = mermaid_code[3:].strip()
        if mermaid_code.endswith("```"):
            mermaid_code = mermaid_code[:-3].strip()
    except Exception:
        logger.exception("LLM生成Mermaid失败")
        mermaid_code = (
            "%%{init: {'theme': 'base', 'themeVariables': { 'primaryColor': '#e8f4fd'}}}%%\n"
            f"graph TD\n    A[{description}] --> B[无法生成详细图表]"
        )

    logger.info("Diagram generated: %s (%d chars)", diagram_type, len(mermaid_code))
    return DocResult(
        filepath="",
        filename="",
        mime_type="text/mermaid",
        content=mermaid_code,
    )


# ── 辅助函数 ──

def _add_text_box(slide, left, top, width, height, text,
                  font_size=None, bold=False, alignment=None):
    """在幻灯片上添加文本框。"""
    from pptx.util import Pt
    from pptx.enum.text import PP_ALIGN
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = text
    if font_size:
        p.font.size = font_size
    p.font.bold = bold
    if alignment is not None:
        p.alignment = alignment


def _parse_slides(raw: str, outline: list[str]) -> dict[str, list[str]]:
    """解析 LLM 输出为 {页标题: [要点列表]}。"""
    import re
    result: dict[str, list[str]] = {}
    current_title = ""
    current_bullets: list[str] = []

    for line in raw.split("\n"):
        stripped = line.strip()
        # 匹配 "## 标题" 或 "## 1. 标题"
        m = re.match(r'^##\s*(?:\d+[\.\、\)]\s*)?(.+)$', stripped)
        if m:
            if current_title and current_bullets:
                result[current_title] = current_bullets
            current_title = m.group(1).strip()
            current_bullets = []
        elif stripped.startswith("- ") or stripped.startswith("* ") or stripped.startswith("• "):
            current_bullets.append(stripped)
        elif stripped and current_title:
            current_bullets.append(f"• {stripped}")

    if current_title and current_bullets:
        result[current_title] = current_bullets

    # 确保所有大纲标题都有内容（仅当完全无匹配时才用通用兜底）
    for s in outline:
        if s not in result:
            result[s] = [f"• [通用知识] {s}相关内容"]

    return result


def _parse_markdown_tables(raw: str) -> list[tuple[str, list[str], list[list[str]]]]:
    """解析 Markdown 中的表格。

    Returns:
        [(table_name, [headers], [[row cells]])]
    """
    import re
    tables = []
    lines = raw.split("\n")
    i = 0
    current_name = ""
    while i < len(lines):
        line = lines[i].strip()
        # 表格名前标注
        m = re.match(r'^##\s*(.+)', line)
        if m:
            current_name = m.group(1).strip()
            i += 1
            continue
        # 检测表头行（以 | 开头）
        if line.startswith("|") and line.endswith("|"):
            headers = [h.strip() for h in line.split("|")[1:-1]]
            # 检查下一行是否是分隔行
            if i + 1 < len(lines):
                sep = lines[i + 1].strip()
                if re.match(r'^[\|\s\-:]+$', sep):
                    i += 2  # 跳过表头和分隔行
                    rows = []
                    while i < len(lines):
                        row_line = lines[i].strip()
                        if row_line.startswith("|") and row_line.endswith("|"):
                            cells = [c.strip() for c in row_line.split("|")[1:-1]]
                            rows.append(cells)
                            i += 1
                        else:
                            break
                    tables.append((current_name, headers, rows))
                    current_name = ""
                    continue
        i += 1

    return tables


def _safe_filename(title: str) -> str:
    """生成安全的文件名。"""
    import re
    safe = re.sub(r'[\\/*?:"<>|]', '', title)
    safe = safe.strip().replace(" ", "_")[:100]
    if not safe:
        safe = "document"
    return safe


def _find_cjk_font() -> str | None:
    """查找系统中可用的 CJK 字体。"""
    candidates = [
        # Windows
        "C:/Windows/Fonts/simsun.ttc",
        "C:/Windows/Fonts/msyh.ttc",
        "C:/Windows/Fonts/simhei.ttf",
        # Linux
        "/usr/share/fonts/truetype/droid/DroidSansFallbackFull.ttf",
        "/usr/share/fonts/opentype/noto/NotoSansCJK-Regular.ttf",
        "/usr/share/fonts/noto-cjk/NotoSansCJK-Regular.ttf",
        # macOS
        "/System/Library/Fonts/PingFang.ttc",
        "/Library/Fonts/Arial Unicode.ttf",
    ]
    for fp in candidates:
        if os.path.isfile(fp):
            logger.info("Found CJK font: %s", fp)
            return fp
    return None
